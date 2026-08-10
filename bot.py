"""Telegram-бот с контекстным AI-ассистентом на Gemini API."""
import asyncio
import json
import logging
import os
import threading
from collections.abc import Iterable
from datetime import timedelta
from io import BytesIO
import time

from dotenv import load_dotenv
from flask import Flask
from google import genai
from google.genai import errors, types
from pptx import Presentation
from wsgiref.simple_server import make_server
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from telegram import Update
from telegram.constants import ChatAction
from telegram.error import TelegramError
from telegram.ext import (
    Application,
    CommandHandler,
    ContextTypes,
    MessageHandler,
    filters,
)

load_dotenv()

# --- ВЕБ-СЕРВЕР ДЛЯ RENDER ---
app = Flask(__name__)

@app.route('/')
def health_check():
    return "Bot is active and running!", 200

def run_web():
    port = int(os.environ.get("PORT", 8080))
    with make_server('0.0.0.0', port, app) as httpd:
        logger.info(f"WSGI server running on port {port}")
        httpd.serve_forever()

# Запускаем Flask в отдельном демоническом потоке
threading.Thread(target=run_web, daemon=True).start()

# --- КОНФИГУРАЦИЯ И КЛЮЧИ ---
# --- КОНФИГУРАЦИЯ И КЛЮЧИ ---
TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
MODEL = os.getenv("GEMINI_MODEL", "gemini-2.0-flash").strip()
IMAGE_MODEL = os.getenv("GEMINI_IMAGE_MODEL", "imagen-3.0-generate-002").strip()
MAX_HISTORY_MESSAGES = 12
TELEGRAM_MESSAGE_LIMIT = 4096
def get_api_keys() -> list[str]:
    load_dotenv(override=True)
    keys = [
        os.environ.get("GEMINI_API_KEY"),
        os.environ.get("GEMINI_API_KEY_2"),
        os.environ.get("GEMINI_API_KEY_3"),
    ]
    return [k for k in keys if k]

API_KEYS = get_api_keys()

IMAGE_COOLDOWN_SECONDS = int(os.environ.get("IMAGE_COOLDOWN_SECONDS", "300"))
# timestamp until which image generation is disabled (epoch seconds)
image_disabled_until = 0.0

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

FALLBACK_TEXT_MODELS = [
    "gemini-2.5-flash-lite",
    "gemini-2.0-flash",
]
FALLBACK_IMAGE_MODELS = [
    "imagen-3.0-generate-002",
    "imagen-2.0-preview",
]


def is_gemini_unavailable_error(error: Exception) -> bool:
    # Prefer structured ClientError checks when available
    try:
        if isinstance(error, errors.ClientError):
            msg = str(error).lower()
            if "429" in msg or "resource_exhausted" in msg or "quota" in msg or "rate limit" in msg:
                return True
            # some ClientError objects may expose a code/status attribute
            code = getattr(error, "code", None) or getattr(error, "status", None) or getattr(error, "status_code", None)
            try:
                if int(code) == 429:
                    return True
            except Exception:
                pass
    except Exception:
        pass

    message = str(error).lower()
    return any(token in message for token in (
        "429",
        "resource_exhausted",
        "quota",
        "rate limit",
        "retry delay",
        "temporarily unavailable",
        "timed out",
        "connection",
        "service unavailable",
    ))


def format_gemini_error_message(error: Exception) -> str:
    if is_gemini_unavailable_error(error):
        return (
            "Сейчас Gemini временно недоступен или достигнут лимит запросов. "
            "Подождите немного и попробуйте ещё раз. Если проблема повторяется, отправьте /status для диагностики."
        )

    message = str(error).lower()
    if "api key" in message or "authentication" in message:
        return "Не удалось подключиться к Gemini. Проверьте API-ключ."

    return "Произошла ошибка при обращении к Gemini. Попробуйте еще раз."


SYSTEM_PROMPT = """
Ты — профессиональный разработчик с многолетним опытом. Объясняешь хорошо и понятно.

ТВОЙ СТИЛЬ:
1. Говоришь профессионально, но доступно — как опытный коллега новичку
2. ВСЕ коды в блоках Markdown с языком: ```python, ```bash, ```javascript и т.д.
3. Структура ответа:
   • Сначала решение (код)
   • Потом объяснение: почему это так, как это работает
   • Если нужно — альтернативы и лучшие практики
4. Объясняй суть, а не только синтаксис. Помогай понять, ЧТО делается и ПОЧЕМУ
5. Примеры из реальной жизни, если помогут пониманию
6. Не используй LaTeX ($ $$) — формулы через Unicode: x², √D, ∑, ≠
"""

# --- ФУНКЦИЯ РОТАЦИИ КЛЮЧЕЙ (БЕЗ БЛОКИРОВКИ EVENT LOOP) ---
def get_models_to_try(model: str, is_image: bool = False) -> list[str]:
    models = [model]
    fallbacks = FALLBACK_IMAGE_MODELS if is_image else FALLBACK_TEXT_MODELS
    for fallback in fallbacks:
        if fallback != model:
            models.append(fallback)
    return models


async def ask_gemini_with_fallback(contents, system_instruction=None, model=None, response_mime_type=None):
    keys = get_api_keys()
    if not keys:
        raise RuntimeError("No GEMINI_API_KEY provided")

    target_model = model or MODEL
    is_image = (
        target_model.startswith("imagen")
        or (target_model.startswith("gemini-3") and "image" in target_model.lower())
    )
    model_candidates = get_models_to_try(target_model, is_image=is_image)
    config = types.GenerateContentConfig(
        system_instruction=system_instruction,
        temperature=0.7,
        max_output_tokens=8192,
        response_mime_type=response_mime_type
    )

    last_error = None
    for model_name in model_candidates:
        for attempt in range(2):
            for i, key in enumerate(keys):
                try:
                    client = genai.Client(api_key=key)
                    response = await client.aio.models.generate_content(
                        model=model_name,
                        contents=contents,
                        config=config
                    )
                    if model_name != target_model:
                        logger.info(f"Успешно отвечено на модели {model_name} после сбоя {target_model}")
                    return response
                except Exception as e:
                    last_error = e
                    logger.warning(f"Ключ #{i+1} ({key[:8]}...) модель {model_name} ошибка на попытке {attempt + 1}: {e}")
                    if not is_gemini_unavailable_error(e):
                        break

            if attempt == 0 and is_gemini_unavailable_error(last_error):
                logger.warning("Временная ошибка Gemini. Повторяем запрос через 2 секунды...")
                await asyncio.sleep(2)
                continue

            break

        if last_error is None or not is_gemini_unavailable_error(last_error):
            break

    if last_error is not None:
        raise last_error

    raise RuntimeError("Gemini request failed")

# --- ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ ---
def split_message(text: str) -> Iterable[str]:
    while text:
        if len(text) <= TELEGRAM_MESSAGE_LIMIT:
            yield text
            return

        split_at = text.rfind("\n", 0, TELEGRAM_MESSAGE_LIMIT)
        if split_at == -1:
            split_at = text.rfind(" ", 0, TELEGRAM_MESSAGE_LIMIT)
        if split_at == -1:
            split_at = TELEGRAM_MESSAGE_LIMIT

        yield text[:split_at]
        text = text[split_at:].lstrip()

# --- ОБРАБОТЧИКИ КОМАНД ---
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    await update.effective_message.reply_text(
        "👋 **Привет! / Welcome!**\n\n"
        "Я AI-ментор по программированию и помощник в учебе. Задай вопрос, отправь код или фото задачи!\n\n"
        "📌 **Команды:**\n"
        "/reset — Очистить историю диалога\n"
        "/slides <тема> — Создать презентацию PPTX\n"
        "/image <описание> — Сгенерировать картинку\n"
        "/remind <минуты> <текст> — Поставить напоминание\n"
        "/help — Помощь и подробности"
    )

async def reset(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    context.user_data.pop("history", None)
    await update.effective_message.reply_text("История диалога очищена. Начнём заново!")

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    await update.effective_message.reply_text(
        "Отправьте текст, фотографию или изображение файлом — я отвечу с учётом "
        "нашего диалога. К изображению можно добавить подпись с вопросом.\n\n"
        "Напоминание: /remind <минуты> <текст>\n"
        "Пример: /remind 30 Выпить воды\n\n"
        "AI-изображение: /image <описание>\n"
        "Презентация: /slides <тема>.\n\n"
        "Используйте /reset, чтобы очистить контекст."
    )


async def status_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Show current API key count, active models and image-generation status."""
    keys = get_api_keys()
    remaining = 0
    if image_disabled_until and image_disabled_until > time.time():
        remaining = int(image_disabled_until - time.time())

    text = (
        f"API keys loaded: {len(keys)}\n"
        f"Text model: {MODEL}\n"
        f"Image model: {IMAGE_MODEL}\n"
    )
    if remaining > 0:
        text += f"Image generation: temporarily disabled (cooldown {remaining}s)"
    else:
        text += "Image generation: enabled"

    await update.effective_message.reply_text(text)

async def send_reminder(context: ContextTypes.DEFAULT_TYPE) -> None:
    job = context.job
    await context.bot.send_message(
        chat_id=job.chat_id,
        text=f"⏰ Напоминание: {job.data['text']}",
    )

async def remind(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    if not context.args or len(context.args) < 2:
        await update.effective_message.reply_text(
            "Формат: /remind <минуты> <текст>\nНапример: /remind 30 Выпить воды"
        )
        return

    try:
        minutes = int(context.args[0])
        if minutes < 1 or minutes > 43_200:
            raise ValueError
    except ValueError:
        await update.effective_message.reply_text("Укажите число минут от 1 до 43200.")
        return

    if context.job_queue is None:
        await update.effective_message.reply_text("Напоминания не настроены.")
        return

    reminder_text = " ".join(context.args[1:])
    context.job_queue.run_once(
        send_reminder,
        when=timedelta(minutes=minutes),
        chat_id=update.effective_chat.id,
        data={"text": reminder_text},
        name=f"reminder-{update.effective_chat.id}",
    )
    await update.effective_message.reply_text(
        f"Хорошо, напомню через {minutes} мин.: {reminder_text}"
    )

def image_from_response(response: object) -> bytes | None:
    for part in getattr(response, "parts", None) or []:
        inline_data = getattr(part, "inline_data", None)
        if inline_data and getattr(inline_data, "data", None):
            return bytes(inline_data.data)
    return None

async def image_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    if not context.args:
        await update.effective_message.reply_text("Напишите, что создать. Например: /image уютный домик в горах")
        return
    prompt = " ".join(context.args)

    global image_disabled_until
    if image_disabled_until and image_disabled_until > time.time():
        remaining = int(image_disabled_until - time.time())
        await update.effective_message.reply_text(
            f"Генерация изображений временно отключена ({remaining}s). Попробуйте позже или используйте текстовый чат."
        )
        return

    try:
        await update.effective_message.reply_text("🎨 Создаю изображение...")
        await update.effective_chat.send_action(ChatAction.UPLOAD_PHOTO)
        response = await ask_gemini_with_fallback(contents=prompt, model=IMAGE_MODEL)
        image_bytes = image_from_response(response)
        if image_bytes is None:
            await update.effective_message.reply_text("Не удалось сгенерировать изображение.")
            return
        await update.effective_message.reply_photo(
            photo=image_bytes,
            caption=f"AI-изображение: {prompt[:900]}",
        )
    except Exception as error:
        logger.exception("Image generation failed: %s", error)
        # If quota/429, disable image generation for a short cooldown to avoid repeated failures
        if is_gemini_unavailable_error(error):
            image_disabled_until = time.time() + IMAGE_COOLDOWN_SECONDS
            await update.effective_message.reply_text(
                f"Генерация изображений временно недоступна из-за лимитов Gemini. Отключено на {IMAGE_COOLDOWN_SECONDS} секунд."
            )
            return

        await update.effective_message.reply_text(format_gemini_error_message(error))

# --- ГЕНЕРАЦИЯ ПРЕЗЕНТАЦИЙ (PPTX) ---
def add_textbox(slide, text: str, left, top, width, height, size: int, color) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = RGBColor(*color)
    return box

def create_presentation(slides: list[dict[str, object]], cover_image: bytes | None) -> BytesIO:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    palettes = [
        ((250, 247, 255), (109, 40, 217), (237, 233, 254)),
        ((240, 249, 255), (3, 105, 161), (224, 242, 254)),
        ((240, 253, 244), (21, 128, 61), (220, 252, 231)),
        ((255, 247, 237), (194, 65, 12), (255, 237, 213)),
    ]

    for index, slide_data in enumerate(slides):
        slide = presentation.slides.add_slide(blank_layout)
        background_color, accent_color, pale_accent = palettes[index % len(palettes)]
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(*background_color)

        blob = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(10.85), Inches(-0.8), Inches(3.0), Inches(3.0)
        )
        blob.fill.solid()
        blob.fill.fore_color.rgb = RGBColor(*pale_accent)
        blob.line.fill.background()

        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.55), Inches(0.78), Inches(0.15)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*accent_color)
        accent_bar.line.fill.background()

        title = str(slide_data.get("title", "Без названия"))[:120]
        add_textbox(
            slide, title, Inches(0.7), Inches(0.85), Inches(10.8), Inches(0.8), 29, (30, 41, 59)
        )

        if index == 0 and cover_image:
            try:
                cover_card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.35), Inches(1.65), Inches(4.2), Inches(4.75)
                )
                cover_card.fill.solid()
                cover_card.fill.fore_color.rgb = RGBColor(*pale_accent)
                cover_card.line.fill.background()
                slide.shapes.add_picture(
                    BytesIO(cover_image), Inches(8.55), Inches(1.85), width=Inches(3.8)
                )
            except (OSError, ValueError):
                logger.warning("Could not insert cover image")

        bullets = slide_data.get("bullets", [])
        if not isinstance(bullets, list):
            bullets = []
        visible_bullets = bullets[:4] or ["Основная идея", "Ключевой вывод"]
        content_width = Inches(7.1) if index == 0 and cover_image else Inches(11.7)
        columns = 1 if len(visible_bullets) < 3 or (index == 0 and cover_image) else 2
        card_width = (content_width - Inches(0.25 * (columns - 1))) / columns

        for bullet_index, bullet in enumerate(visible_bullets):
            column = bullet_index % columns
            row = bullet_index // columns
            card_height = Inches(1.38)
            card_left = Inches(0.7) + column * (card_width + Inches(0.25))
            card_top = Inches(1.85) + row * Inches(1.62)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(*pale_accent)
            number = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, card_left + Inches(0.22), card_top + Inches(0.25), Inches(0.48), Inches(0.48)
            )
            number.fill.solid()
            number.fill.fore_color.rgb = RGBColor(*accent_color)
            number.line.fill.background()
            number_text = add_textbox(
                slide, str(bullet_index + 1), card_left + Inches(0.22), card_top + Inches(0.33), Inches(0.48), Inches(0.2), 10, (255, 255, 255)
            )
            number_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            card_text = add_textbox(
                slide, str(bullet)[:150], card_left + Inches(0.88), card_top + Inches(0.2), card_width - Inches(1.08), Inches(0.95), 16, (51, 65, 85)
            )
            card_text.text_frame.word_wrap = True

        footer = add_textbox(
            slide, f"{index + 1} / {len(slides)}", Inches(11.7), Inches(6.85), Inches(0.8), Inches(0.3), 10, (100, 116, 139)
        )
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    file = BytesIO()
    presentation.save(file)
    file.seek(0)
    return file

async def slides_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    command_args = context.args
    caption = update.effective_message.caption or ""
    if not command_args and caption.startswith("/slides"):
        command_args = caption.split()[1:]
    if not command_args:
        await update.effective_message.reply_text("Укажите тему: /slides Искусственный интеллект")
        return

    message = update.effective_message
    image = message.photo[-1] if message.photo else message.document
    image_bytes = None
    image_part = None
    if image:
        try:
            file = await image.get_file()
            image_bytes = bytes(await file.download_as_bytearray())
            image_part = types.Part.from_bytes(
                data=image_bytes,
                mime_type=getattr(image, "mime_type", None) or "image/jpeg",
            )
        except (OSError, ValueError, TelegramError):
            await message.reply_text("Не удалось скачать картинку для презентации.")
            return

    topic = " ".join(command_args)
    prompt = (
        f"Составь план презентации на русском языке по теме: {topic}. "
        "Верни только корректный JSON без Markdown: объект с ключом slides. "
        "slides — массив из 5–7 объектов; у каждого title (короткий заголовок) "
        "и bullets (массив из 3–5 коротких тезисов)."
    )
    contents = [types.Part.from_text(text=prompt)]
    if image_part:
        contents.append(image_part)

    try:
        await update.effective_chat.send_action(ChatAction.TYPING)
        response = await ask_gemini_with_fallback(
            contents=contents,
            response_mime_type="application/json"
        )
        outline = json.loads((response.text or "").strip())
        slides = outline.get("slides", [])
        if not isinstance(slides, list) or not 2 <= len(slides) <= 10:
            raise ValueError("invalid slide outline")
        document = create_presentation(slides, image_bytes)
    except Exception:
        logger.exception("Presentation generation failed")
        await message.reply_text("Не удалось создать презентацию. Попробуйте еще раз.")
        return

    await update.effective_chat.send_action(ChatAction.UPLOAD_DOCUMENT)
    await message.reply_document(
        document=document,
        filename="presentation.pptx",
        caption=f"Презентация по теме «{topic[:120]}» готова!"
    )

# --- ОСНОВНОЙ ЧАТ С ИСТОРИЕЙ ---
async def chat(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = update.effective_message
    user_text = (message.text or message.caption or "").strip()
    if message.caption and message.caption.startswith("/slides"):
        await slides_command(update, context)
        return
    photo = message.photo[-1] if message.photo else None
    image = photo or message.document
    if not user_text and image is None:
        return

    history: list[dict[str, str]] = context.user_data.get("history", [])
    contents = []

    # Добавляем историю сообщений
    for item in history:
        contents.append({
            "role": item["role"],
            "parts": [{"text": item["content"]}]
        })

    # Добавляем текущий запрос
    user_parts = []
    if image:
        try:
            file = await image.get_file()
            image_bytes = await file.download_as_bytearray()
            user_parts.append(
                types.Part.from_bytes(
                    data=bytes(image_bytes),
                    mime_type=getattr(image, "mime_type", None) or "image/jpeg",
                )
            )
        except (OSError, ValueError, TelegramError) as e:
            logger.exception("Telegram photo download failed: %s", e)
            await message.reply_text("Не удалось скачать изображение.")
            return

    if user_text or not image:
        user_parts.append(types.Part.from_text(text=user_text or "Опиши и проанализируй это изображение."))

    contents.append({
        "role": "user",
        "parts": user_parts
    })

    try:
        await update.effective_chat.send_action(
            ChatAction.UPLOAD_PHOTO if image else ChatAction.TYPING
        )

        response = await ask_gemini_with_fallback(
            contents=contents,
            system_instruction=SYSTEM_PROMPT
        )

        reply = (response.text or "").strip() or "Не удалось сформировать ответ. Попробуйте переформулировать."

    except Exception as error:
        logger.error("Ошибка в функции chat: %s", error, exc_info=True)
        await update.effective_message.reply_text(format_gemini_error_message(error))
        return

    # Сохраняем успешный ответ в историю
    updated_history = [
        *history,
        {
            "role": "user",
            "content": user_text or "[Отправлено изображение]",
        },
        {"role": "model", "content": reply},
    ]
    context.user_data["history"] = updated_history[-MAX_HISTORY_MESSAGES:]

    # Отправляем ответ частями
    for part in split_message(reply):
        try:
            await update.effective_message.reply_text(part, parse_mode="Markdown")
        except Exception:
            await update.effective_message.reply_text(part)

async def error_handler(update: object, context: ContextTypes.DEFAULT_TYPE) -> None:
    logger.exception("Unhandled exception while processing update", exc_info=context.error)

# --- ТОЧКА ВХОДА ---
def main() -> None:
    if not TELEGRAM_BOT_TOKEN:
        raise RuntimeError("TELEGRAM_BOT_TOKEN is missing from .env")
    if not API_KEYS:
        raise RuntimeError("No GEMINI_API_KEY provided")

    application = Application.builder().token(TELEGRAM_BOT_TOKEN).build()
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("reset", reset))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CommandHandler("remind", remind))
    application.add_handler(CommandHandler("image", image_command))
    application.add_handler(CommandHandler("status", status_command))
    application.add_handler(CommandHandler("slides", slides_command))
    application.add_handler(
        MessageHandler(
            (filters.TEXT & ~filters.COMMAND) | filters.PHOTO | filters.Document.IMAGE,
            chat,
        )
    )
    application.add_error_handler(error_handler)

    logger.info("Bot starting...")
    application.run_polling(drop_pending_updates=True)

if __name__ == "__main__":
    main()